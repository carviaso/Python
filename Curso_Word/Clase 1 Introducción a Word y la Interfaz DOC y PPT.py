from pptx import Presentation

# Crear presentación PowerPoint
ppt = Presentation()

# Diapositiva 1 - Título
slide_1 = ppt.slides.add_slide(ppt.slide_layouts[0])
title = slide_1.shapes.title
title.text = "Curso de Microsoft Word Básico - Clase 1"
subtitle = slide_1.placeholders[1]
subtitle.text = "Introducción a Word y la Interfaz"

# Diapositiva 2 - Objetivos de la clase
slide_2 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Objetivos de la Clase"
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• Conocer qué es Word y sus usos.\n• Explorar las partes principales de la interfaz.\n• Aprender a crear, abrir y guardar documentos."

# Diapositiva 3 - ¿Qué es Word y para qué se usa?
slide_3 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_3.shapes.title
title.text = "¿Qué es Word y para qué se usa?"
content = slide_3.shapes.placeholders[1].text_frame
content.text = "• Word es un procesador de texto utilizado para crear y editar documentos.\n• Usos: escribir cartas, informes, currículums, etc."

# Diapositiva 4 - Partes principales de la interfaz
slide_4 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Partes principales de la interfaz"
content = slide_4.shapes.placeholders[1].text_frame
content.text = "• Barra de herramientas\n• Cinta de opciones\n• Área de trabajo\n• Barra de estado"

# Diapositiva 5 - Crear, abrir y guardar documentos
slide_5 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Crear, abrir y guardar documentos"
content = slide_5.shapes.placeholders[1].text_frame
content.text = "• Para crear un documento, selecciona 'Nuevo' o 'Nuevo en blanco'.\n• Para abrir, selecciona 'Abrir'.\n• Guardar documento: Ctrl + S"

# Diapositiva 6 - Práctica Guiada
slide_6 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Práctica Guiada"
content = slide_6.shapes.placeholders[1].text_frame
content.text = "• Explorar la interfaz de Word.\n• Crear un nuevo documento y guardarlo en diferentes formatos."

# Diapositiva 7 - Ejercicios y Repaso
slide_7 = ppt.slides.add_slide(ppt.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Ejercicios y Repaso"
content = slide_7.shapes.placeholders[1].text_frame
content.text = "• Crear un documento con texto de ejemplo.\n• Guardar el documento en formato .docx y .pdf"

# Guardar presentación
pptx_filename = "D:/Clase_1_Curso_Word_Basico.pptx"
ppt.save(pptx_filename)

# Crear archivo Word para ejercicio
from docx import Document

doc = Document()
doc.add_heading("Ejercicio: Crear un Documento y Guardarlo", 0)

doc.add_paragraph(
    "Instrucciones:\n\n"
    "1. Crea un documento nuevo en Word.\n"
    "2. Escribe un texto de ejemplo, por ejemplo, tu nombre y una breve introducción sobre ti.\n"
    "3. Guarda el documento en formato .docx.\n"
    "4. Vuelve a guardar el documento en formato .pdf.\n"
    "5. Envía ambos archivos para su revisión.\n\n"
    "Formato esperado: Un archivo .docx y un archivo .pdf."
)

# Guardar documento Word
docx_filename = "d:/Clase_1_Curso_Word_Basico_Ejercicio.docx"
doc.save(docx_filename)

pptx_filename, docx_filename
