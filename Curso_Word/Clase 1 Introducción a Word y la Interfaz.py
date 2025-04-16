from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Crear la presentación
prs = Presentation()

# Estilos básicos
title_font_size = Pt(44)
subtitle_font_size = Pt(32)
content_font_size = Pt(20)

def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

def add_content_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]

    title_shape.text = title
    tf = body_shape.text_frame
    tf.clear()  # limpiar contenido existente

    for i, point in enumerate(bullet_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = content_font_size

# Crear las diapositivas para la Clase 1
add_title_slide("Clase 1: Introducción a Word y la Interfaz", "Curso de Word Básico")

add_content_slide("¿Qué es Word y para qué se usa?", [
    "Microsoft Word es un procesador de textos.",
    "Permite crear, editar y dar formato a documentos de texto.",
    "Se usa para informes, cartas, currículos y más."
])

add_content_slide("Partes principales de la interfaz", [
    "Barra de herramientas de acceso rápido",
    "Cinta de opciones con pestañas",
    "Área de trabajo (hoja en blanco)",
    "Reglas, barra de estado y de desplazamiento"
])

add_content_slide("Crear, abrir y guardar documentos", [
    "Crear: iniciar un nuevo documento en blanco o desde plantilla.",
    "Abrir: buscar y cargar un archivo existente.",
    "Guardar: almacenar el archivo en el equipo o en la nube.",
    "Guardar como PDF: útil para compartir sin que se modifique."
])

# Guardar el archivo
pptx_path = "d:/Clase_1_Word_Introduccion.pptx"
prs.save(pptx_path)

pptx_path
