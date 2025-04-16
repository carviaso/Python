# -*- coding: utf-8 -*-
"""
Created on Tue Apr 15 08:26:15 2025
Clase 6, 7, 8, 9, 10

@author: chatGpt
"""

from pptx import Presentation
from docx import Document

# Función para crear una presentación PowerPoint
def crear_ppt(nombre_archivo, titulo, subtitulo, contenido_diapositivas):
    ppt = Presentation()
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = subtitulo

    for titulo_slide, contenido in contenido_diapositivas:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        slide.shapes.title.text = titulo_slide
        slide.placeholders[1].text = contenido

    ruta = f"d:/{nombre_archivo}.pptx"
    ppt.save(ruta)
    return ruta

# Función para crear un documento Word
def crear_doc(nombre_archivo, titulo, ejercicios):
    doc = Document()
    doc.add_heading(titulo, 0)
    for i, (titulo_ej, instrucciones) in enumerate(ejercicios, 1):
        doc.add_paragraph(f"Ejercicio {i}: {titulo_ej}")
        doc.add_paragraph(instrucciones)
    ruta = f"d:/{nombre_archivo}.docx"
    doc.save(ruta)
    return ruta

# Clase 6
ppt6 = crear_ppt(
    "Clase_06_Imagenes_Graficos_Word",
    "Clase 6: Imágenes y Gráficos",
    "Curso de Microsoft Word Básico\nDuración: 4 horas",
    [
        ("Objetivos de la Clase", "- Insertar imágenes, ajustar texto, insertar gráficos y formas."),
        ("Insertar Imágenes", "- Insertar desde archivo o en línea.\n- Ajustar tamaño y posición."),
        ("Ajuste de Texto", "- Usar opciones de ajuste de texto: cuadrado, estrecho, detrás del texto."),
        ("Insertar Gráficos y Formas", "- Insertar formas básicas, SmartArt y gráficos.\n- Modificar colores y estilos."),
        ("Práctica Guiada", "- Insertar una imagen y aplicar diferentes tipos de ajuste.\n- Crear un gráfico con formas.")
    ]
)

doc6 = crear_doc(
    "Practica_Clase_06_Word",
    "Clase 6: Práctica - Imágenes y Gráficos",
    [
        ("Insertar y Ajustar Imágenes", "1. Inserte una imagen.\n2. Pruebe al menos tres tipos de ajuste de texto."),
        ("Insertar Formas y Gráficos", "1. Inserte 3 formas y combínelas.\n2. Cree un SmartArt con 3 elementos.\n3. Inserte un gráfico de barras con datos de ventas ficticios.")
    ]
)

# Clase 7
ppt7 = crear_ppt(
    "Clase_07_Herramientas_Revisión_Word",
    "Clase 7: Herramientas de Revisión",
    "Curso de Microsoft Word Básico\nDuración: 4 horas",
    [
        ("Objetivos de la Clase", "- Usar corrector ortográfico, comentarios y control de cambios."),
        ("Corrector Ortográfico", "- Activar corrector y revisar sugerencias."),
        ("Insertar Comentarios", "- Seleccionar texto y agregar comentario."),
        ("Control de Cambios", "- Activar y ver ediciones.\n- Aceptar o rechazar cambios."),
        ("Práctica Guiada", "- Revisar un texto con errores.\n- Insertar comentarios y usar control de cambios.")
    ]
)

doc7 = crear_doc(
    "Practica_Clase_07_Word",
    "Clase 7: Práctica - Herramientas de Revisión",
    [
        ("Revisión Ortográfica", "Corrija un texto con errores ortográficos."),
        ("Comentarios", "Agregue 3 comentarios en un párrafo de texto."),
        ("Control de Cambios", "Active control de cambios y edite un párrafo.\nRevise los cambios y acepte algunos.")
    ]
)

# Clase 8
ppt8 = crear_ppt(
    "Clase_08_Impresion_Exportacion_Word",
    "Clase 8: Impresión y Exportación",
    "Curso de Microsoft Word Básico\nDuración: 4 horas",
    [
        ("Objetivos de la Clase", "- Configurar impresión y guardar en diferentes formatos."),
        ("Configuración de Impresión", "- Usar vista previa y ajustar configuración: orientación, tamaño, páginas."),
        ("Exportar Documentos", "- Guardar como PDF, DOCX y otros."),
        ("Práctica Guiada", "- Configurar un documento para impresión a doble cara.\n- Exportar a PDF.")
    ]
)

doc8 = crear_doc(
    "Practica_Clase_08_Word",
    "Clase 8: Práctica - Impresión y Exportación",
    [
        ("Configurar impresión", "1. Configure un documento para impresión doble cara.\n2. Ajuste márgenes personalizados."),
        ("Exportar documento", "1. Guarde el documento como PDF.\n2. Guarde como archivo de solo lectura.")
    ]
)

# Clase 9
ppt9 = crear_ppt(
    "Clase_09_Tablas_Contenido_Estilos",
    "Clase 9: Tabla de Contenido y Estilos",
    "Curso de Microsoft Word Básico\nDuración: 4 horas",
    [
        ("Objetivos de la Clase", "- Crear tabla de contenido automática.\n- Usar estilos."),
        ("Aplicar Estilos", "- Usar estilos predefinidos para títulos y subtítulos."),
        ("Insertar Tabla de Contenido", "- Usar 'Referencias > Tabla de contenido'.\n- Actualizar tabla automáticamente."),
        ("Práctica Guiada", "- Crear un documento con títulos y subtítulos.\n- Insertar y actualizar tabla de contenido.")
    ]
)

doc9 = crear_doc(
    "Practica_Clase_09_Word",
    "Clase 9: Práctica - Tabla de Contenido y Estilos",
    [
        ("Aplicar Estilos", "1. Agregue encabezados (Título 1, Título 2) a un documento."),
        ("Insertar Tabla de Contenido", "1. Inserte tabla de contenido automática.\n2. Modifique el texto y actualice la tabla.")
    ]
)

# Clase 10
ppt10 = crear_ppt(
    "Clase_10_Proyecto_Final_Repaso",
    "Clase 10: Proyecto Final y Repaso",
    "Curso de Microsoft Word Básico\nDuración: 4 horas",
    [
        ("Objetivos de la Clase", "- Reforzar conocimientos y realizar un proyecto final."),
        ("Proyecto Final", "- Crear un documento formal completo.\n- Incluir portada, índice, texto, imágenes y tabla."),
        ("Repaso General", "- Preguntas tipo quiz.\n- Resolución de dudas."),
        ("Entrega y Evaluación", "- Entregar el documento y revisión del instructor.")
    ]
)

doc10 = crear_doc(
    "Practica_Clase_10_Word",
    "Clase 10: Práctica - Proyecto Final y Repaso",
    [
        ("Proyecto Final", "Cree un documento completo sobre un tema libre que incluya:\n- Portada\n- Índice automático\n- Texto con estilos\n- Imágenes y tabla"),
        ("Repaso", "Responda un breve quiz de repaso.\nPreguntas de opción múltiple sobre lo aprendido.")
    ]
)

ppt6, doc6, ppt7, doc7, ppt8, doc8, ppt9, doc9, ppt10, doc10
