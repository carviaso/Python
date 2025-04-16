from pptx import Presentation

# Crear presentación PowerPoint para Clase 2
ppt2 = Presentation()

# Diapositiva 1 - Título
slide_1 = ppt2.slides.add_slide(ppt2.slide_layouts[0])
title = slide_1.shapes.title
title.text = "Curso de Microsoft Word Básico - Clase 2"
subtitle = slide_1.placeholders[1]
subtitle.text = "Edición y Formato de Texto"

# Diapositiva 2 - Objetivos de la clase
slide_2 = ppt2.slides.add_slide(ppt2.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Objetivos de la Clase"
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• Conocer cómo editar texto en Word.\n• Aprender a dar formato al texto.\n• Aplicar efectos de texto como negrita, cursiva y subrayado."

# Diapositiva 3 - Tipo de fuente, tamaño y color
slide_3 = ppt2.slides.add_slide(ppt2.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Tipo de fuente, tamaño y color"
content = slide_3.shapes.placeholders[1].text_frame
content.text = "• Cambiar el tipo de fuente (Ej. Arial, Times New Roman).\n• Ajustar el tamaño de la fuente.\n• Aplicar color al texto."

# Diapositiva 4 - Negrita, cursiva, subrayado y efectos de texto
slide_4 = ppt2.slides.add_slide(ppt2.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Negrita, Cursiva, Subrayado y Efectos de Texto"
content = slide_4.shapes.placeholders[1].text_frame
content.text = "• Negrita (Ctrl + B)\n• Cursiva (Ctrl + I)\n• Subrayado (Ctrl + U)\n• Efectos de texto como sombras y reflejos"

# Diapositiva 5 - Práctica Guiada
slide_5 = ppt2.slides.add_slide(ppt2.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Práctica Guiada"
content = slide_5.shapes.placeholders[1].text_frame
content.text = "• Cambiar el tipo de fuente y tamaño en un documento.\n• Aplicar negrita, cursiva y subrayado en el texto."

# Diapositiva 6 - Ejercicios y Repaso
slide_6 = ppt2.slides.add_slide(ppt2.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Ejercicios y Repaso"
content = slide_6.shapes.placeholders[1].text_frame
content.text = "• Crear un documento con texto utilizando distintos estilos de fuente.\n• Aplicar formato a los párrafos y experimentar con efectos de texto."

# Guardar presentación
pptx_filename_2 = "D:/Clase_2_Curso_Word_Basico.pptx"
ppt2.save(pptx_filename_2)

# Crear archivo Word para ejercicio
from docx import Document

doc2 = Document()
doc2.add_heading("Ejercicio: Edición y Formato de Texto", 0)

doc2.add_paragraph(
    "Instrucciones:\n\n"
    "1. Crea un documento en Word.\n"
    "2. Escribe un párrafo con texto de ejemplo.\n"
    "3. Aplica cambios de formato como negrita, cursiva, subrayado, y cambia el tamaño y color del texto.\n"
    "4. Guarda el documento y envíalo para su revisión.\n\n"
    "Formato esperado: Un archivo con texto formateado de acuerdo a las instrucciones."
)

# Guardar documento Word
docx_filename_2 = "D:/Clase_2_Curso_Word_Basico_Ejercicio.docx"
doc2.save(docx_filename_2)

pptx_filename_2, docx_filename_2
