from pptx import Presentation

# Crear presentación PowerPoint para Clase 3
ppt3 = Presentation()

# Diapositiva 1 - Título
slide_1 = ppt3.slides.add_slide(ppt3.slide_layouts[0])
title = slide_1.shapes.title
title.text = "Curso de Microsoft Word Básico - Clase 3"
subtitle = slide_1.placeholders[1]
subtitle.text = "Formato de Párrafos y Espaciado"

# Diapositiva 2 - Objetivos de la clase
slide_2 = ppt3.slides.add_slide(ppt3.slide_layouts[1])
title = slide_2.shapes.title
title.text = "Objetivos de la Clase"
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• Conocer las opciones de alineación y espaciado en los párrafos.\n• Aprender a aplicar sangrías y tabulaciones.\n• Formatear párrafos de acuerdo a las necesidades."

# Diapositiva 3 - Alineación (izquierda, derecha, centrado, justificado)
slide_3 = ppt3.slides.add_slide(ppt3.slide_layouts[1])
title = slide_3.shapes.title
title.text = "Alineación (izquierda, derecha, centrado, justificado)"
content = slide_3.shapes.placeholders[1].text_frame
content.text = "• Alineación izquierda: texto alineado a la izquierda.\n• Alineación derecha: texto alineado a la derecha.\n• Alineación centrada: texto centrado.\n• Alineación justificada: texto alineado de manera uniforme."

# Diapositiva 4 - Interlineado y espaciado entre párrafos
slide_4 = ppt3.slides.add_slide(ppt3.slide_layouts[1])
title = slide_4.shapes.title
title.text = "Interlineado y espaciado entre párrafos"
content = slide_4.shapes.placeholders[1].text_frame
content.text = "• Interlineado: espacio entre líneas de un párrafo.\n• Espaciado entre párrafos: espacio adicional entre párrafos.\n• Ajustar para mejorar la legibilidad del documento."

# Diapositiva 5 - Sangrías y tabulaciones
slide_5 = ppt3.slides.add_slide(ppt3.slide_layouts[1])
title = slide_5.shapes.title
title.text = "Sangrías y tabulaciones"
content = slide_5.shapes.placeholders[1].text_frame
content.text = "• Sangrías: mover el texto hacia adentro desde el margen.\n• Tabulaciones: mover el texto a una posición específica dentro del párrafo."

# Diapositiva 6 - Práctica Guiada
slide_6 = ppt3.slides.add_slide(ppt3.slide_layouts[1])
title = slide_6.shapes.title
title.text = "Práctica Guiada"
content = slide_6.shapes.placeholders[1].text_frame
content.text = "• Aplicar diferentes alineaciones y sangrías en un documento.\n• Ajustar interlineado y espaciado entre párrafos."

# Diapositiva 7 - Ejercicios y Repaso
slide_7 = ppt3.slides.add_slide(ppt3.slide_layouts[1])
title = slide_7.shapes.title
title.text = "Ejercicios y Repaso"
content = slide_7.shapes.placeholders[1].text_frame
content.text = "• Crear un documento con diferentes estilos de párrafo.\n• Aplicar alineación, interlineado, y espaciado entre párrafos."

# Guardar presentación
pptx_filename_3 = "D:/Clase_3_Curso_Word_Basico.pptx"
ppt3.save(pptx_filename_3)

# Crear archivo Word para ejercicio
from docx import Document

doc3 = Document()
doc3.add_heading("Ejercicio: Formato de Párrafos y Espaciado", 0)

doc3.add_paragraph(
    "Instrucciones:\n\n"
    "1. Crea un documento en Word.\n"
    "2. Escribe varios párrafos de texto.\n"
    "3. Aplica diferentes alineaciones (izquierda, derecha, centrado, justificado).\n"
    "4. Ajusta el interlineado y espaciado entre los párrafos.\n"
    "5. Aplica sangrías y tabulaciones en los párrafos.\n\n"
    "Formato esperado: Un archivo con párrafos correctamente alineados y formateados."
)

# Guardar documento Word
docx_filename_3 = "D:/Clase_3_Curso_Word_Basico_Ejercicio.docx"
doc3.save(docx_filename_3)

pptx_filename_3, docx_filename_3
